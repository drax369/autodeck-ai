import os
import pandas as pd
from pptx import Presentation

DATA_DIR = "data/auto-deck-ai-systems-challenge/Dataset/Dataset"
TRAIN_CSV = "data/auto-deck-ai-systems-challenge/train.csv"
SUBMISSION_CSV = "data/auto-deck-ai-systems-challenge/sample_submission.csv"

def extract_features(filepath, file_extension):
    try:
        if file_extension.lower() in ['.pptx', '.ppsx', '.pptm']:
            prs = Presentation(filepath)
            slide_count = len(prs.slides)
            word_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        word_count += len(shape.text.split())
            return slide_count, word_count, True
        else:
            size_kb = os.path.getsize(filepath) / 1024
            estimated_slides = max(1, int(size_kb / 100))
            size_per_slide = size_kb / max(estimated_slides, 1)
            if size_per_slide < 30:
                words_per_slide = 15
            elif size_per_slide < 80:
                words_per_slide = 45
            else:
                words_per_slide = 90
            estimated_words = estimated_slides * words_per_slide
            return estimated_slides, estimated_words, False
    except Exception:
        size_kb = os.path.getsize(filepath) / 1024
        estimated_slides = max(1, int(size_kb / 100))
        size_per_slide = size_kb / max(estimated_slides, 1)
        if size_per_slide < 30:
            words_per_slide = 15
        elif size_per_slide < 80:
            words_per_slide = 45
        else:
            words_per_slide = 90
        estimated_words = estimated_slides * words_per_slide
        return estimated_slides, estimated_words, False


def classify(slide_count, word_count, file_extension,
             title, company, creation_date, last_modified, exact=True):

    # Complexity
    if slide_count <= 8:
        complexity = "Low"
    elif slide_count <= 25:
        complexity = "Medium"
    else:
        complexity = "High"

    # Density
    avg = word_count / max(slide_count, 1)
    if avg <= 25:
        density = "Light"
    elif avg <= 60:
        density = "Medium"
    else:
        density = "Dense"

    # Category
    ext = file_extension.lower()
    if ext == '.pptx':
        category = "Modern"
    elif ext == '.ppt':
        category = "Legacy"
    elif ext in ['.pps', '.ppsx']:
        category = "Slideshow"
    elif ext in ['.pot', '.potx']:
        category = "Template"
    else:
        category = "Modern"

    # Quality score
    score = 0
    if title and title != '-':
        score += 25
    if company and company != '-':
        score += 25
    if creation_date and creation_date != '-':
        score += 25
    if last_modified and last_modified != '-':
        score += 25

    # Confidence
    confidence = 0.95 if exact else 0.55

    return complexity, density, category, score, round(confidence, 2)


# Load data
print("Loading data...")
df = pd.read_csv(TRAIN_CSV)
submission = pd.read_csv(SUBMISSION_CSV)

test_df = df[df['split'] == 'test'].copy()
print(f"Test files to predict: {len(test_df)}")

results = []
exact_count = 0
estimated_count = 0

for idx, row in test_df.iterrows():
    file_path = os.path.join(DATA_DIR, "test", row['filename'])

    slide_count, word_count, exact = extract_features(file_path, row['file_extension'])

    if exact:
        exact_count += 1
    else:
        estimated_count += 1

    complexity, density, category, quality, confidence = classify(
        slide_count, word_count,
        row['file_extension'],
        row['title'], row['company'],
        row['creation_date'], row['last_modified'],
        exact
    )

    results.append({
        'file_id': row['file_id'],
        'filename': row['filename'],
        'predicted_complexity': complexity,
        'predicted_density': density,
        'predicted_file_category': category,
        'predicted_quality_score': quality,
        'confidence': confidence
    })

    if len(results) % 10 == 0:
        print(f"Processed {len(results)}/{len(test_df)} files...")

output_df = pd.DataFrame(results)
output_df.to_csv('submission.csv', index=False)

print(f"\nDone! submission.csv saved.")
print(f"Exact extractions: {exact_count} | Estimated: {estimated_count}")
print(f"\nPrediction distribution:")
print("Complexity:", output_df['predicted_complexity'].value_counts().to_dict())
print("Density:   ", output_df['predicted_density'].value_counts().to_dict())
print("Category:  ", output_df['predicted_file_category'].value_counts().to_dict())
print("Quality:   ", output_df['predicted_quality_score'].value_counts().to_dict())
print("\nSample output:")
print(output_df.head(5).to_string())