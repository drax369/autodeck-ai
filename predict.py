import os
import pandas as pd
import numpy as np
import pickle
from pptx import Presentation

DATA_DIR = "data/auto-deck-ai-systems-challenge/Dataset/Dataset"
TRAIN_CSV = "data/auto-deck-ai-systems-challenge/train.csv"

# Load models
with open('models.pkl', 'rb') as f:
    models = pickle.load(f)

clf_complexity = models['clf_complexity']
clf_density = models['clf_density']
le_complexity = models['le_complexity']
le_density = models['le_density']

def extract_features(filepath, file_extension):
    try:
        if file_extension.lower() in ['.pptx', '.ppsx', '.pptm']:
            prs = Presentation(filepath)
            slide_count = len(prs.slides)
            word_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        word_count += len(shape.text.split())
            return slide_count, word_count, True
        else:
            size_kb = os.path.getsize(filepath) / 1024
            if size_kb < 297:
                estimated_slides = max(1, int(size_kb / 30))
                words_per_slide = 20
            elif size_kb < 944:
                estimated_slides = max(1, int(size_kb / 80))
                words_per_slide = 45
            elif size_kb < 3136:
                estimated_slides = max(1, int(size_kb / 120))
                words_per_slide = 65
            else:
                estimated_slides = max(1, int(size_kb / 200))
                words_per_slide = 90
            return estimated_slides, estimated_slides * words_per_slide, False
    except Exception:
        size_kb = os.path.getsize(filepath) / 1024
        estimated_slides = max(1, int(size_kb / 100))
        return estimated_slides, estimated_slides * 45, False

def get_category(ext):
    ext = ext.lower()
    if ext == '.pptx': return "Modern"
    elif ext == '.ppt': return "Legacy"
    elif ext in ['.pps', '.ppsx']: return "Slideshow"
    elif ext in ['.pot', '.potx']: return "Template"
    else: return "Modern"

# Load data
df = pd.read_csv(TRAIN_CSV)
test_df = df.copy()
print(f"Predicting {len(test_df)} test files...")

results = []

for i, (idx, row) in enumerate(test_df.iterrows()):
    file_path = os.path.join(DATA_DIR, row['split'], row['filename'])

    slide_count, word_count, exact = extract_features(file_path, row['file_extension'])
    avg_words = word_count / max(slide_count, 1)
    size_kb = row['file_size_bytes'] / 1024
    size_mb = row['file_size_mb']

    X = [[size_kb, size_mb, slide_count, word_count, avg_words]]

    complexity = le_complexity.inverse_transform(clf_complexity.predict(X))[0]
    density = le_density.inverse_transform(clf_density.predict(X))[0]
    category = get_category(row['file_extension'])

    complexity_proba = clf_complexity.predict_proba(X).max()
    density_proba = clf_density.predict_proba(X).max()
    confidence = round((complexity_proba + density_proba) / 2, 2)

    results.append({
        'file_id': row['file_id'],
        'filename': row['filename'],
        'predicted_complexity': complexity,
        'predicted_density': density,
        'predicted_file_category': category,
        'predicted_quality_score': 0,
        'confidence': confidence
    })

    if (i + 1) % 100 == 0:
        print(f"  {i+1}/{len(test_df)} done...")

output_df = pd.DataFrame(results)
output_df.to_csv('submission.csv', index=False)

print(f"\nDone! submission.csv saved.")
print(f"\nPrediction distribution:")
print("Complexity:", output_df['predicted_complexity'].value_counts().to_dict())
print("Density:   ", output_df['predicted_density'].value_counts().to_dict())
print("Category:  ", output_df['predicted_file_category'].value_counts().to_dict())
print("Confidence:", output_df['confidence'].describe().round(2).to_dict())
print("\nSample output:")
print(output_df.head(5).to_string())