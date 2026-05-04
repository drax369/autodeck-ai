import os
import pandas as pd
from pptx import Presentation

DATA_DIR = "data/auto-deck-ai-systems-challenge/Dataset/Dataset"
TRAIN_CSV = "data/auto-deck-ai-systems-challenge/train.csv"

def extract_features(filepath, file_extension):
    try:
        if file_extension.lower() in ['.pptx', '.ppsx', '.pptm']:
            prs = Presentation(filepath)
            slide_count = len(prs.slides)
            word_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        word_count += len(shape.text.split())
            return slide_count, word_count
        else:
            size_kb = os.path.getsize(filepath) / 1024
            if size_kb < 297:
                estimated_slides = max(1, int(size_kb / 30))
                words_per_slide = 20
            elif size_kb < 944:
                estimated_slides = max(1, int(size_kb / 80))
                words_per_slide = 45
            elif size_kb < 3136:
                estimated_slides = max(1, int(size_kb / 120))
                words_per_slide = 65
            else:
                estimated_slides = max(1, int(size_kb / 200))
                words_per_slide = 90
            return estimated_slides, estimated_slides * words_per_slide
    except Exception:
        size_kb = os.path.getsize(filepath) / 1024
        estimated_slides = max(1, int(size_kb / 100))
        return estimated_slides, estimated_slides * 45

df = pd.read_csv(TRAIN_CSV)
train_df = df[df['split'] == 'train'].copy()
print(f"Extracting features from {len(train_df)} training files...")

rows = []
for i, (idx, row) in enumerate(train_df.iterrows()):
    filepath = os.path.join(DATA_DIR, "train", row['filename'])
    slide_count, word_count = extract_features(filepath, row['file_extension'])
    size_kb = row['file_size_bytes'] / 1024
    avg_words = word_count / max(slide_count, 1)

    rows.append({
        'file_id': row['file_id'],
        'filename': row['filename'],
        'file_size_kb': size_kb,
        'file_size_mb': row['file_size_mb'],
        'slide_count': slide_count,
        'word_count': word_count,
        'avg_words_per_slide': avg_words,
        'file_extension': row['file_extension'],
    })

    if (i + 1) % 50 == 0:
        print(f"  {i+1}/{len(train_df)} done...")

features_df = pd.DataFrame(rows)
features_df.to_csv('train_features.csv', index=False)
print(f"\nDone! train_features.csv saved.")
print(features_df.describe().round(2))