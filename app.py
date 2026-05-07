import os
import pickle
import time
import gradio as gr
from pptx import Presentation
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from fpdf import FPDF
import tempfile

# Load models
with open('models.pkl', 'rb') as f:
    models = pickle.load(f)

clf_complexity = models['clf_complexity']
clf_density = models['clf_density']
le_complexity = models['le_complexity']
le_density = models['le_density']

# History storage
analysis_history = []

def extract_features(filepath, file_extension):
    try:
        if file_extension.lower() in ['.pptx', '.ppsx', '.pptm']:
            prs = Presentation(filepath)
            slide_count = len(prs.slides)
            word_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        word_count += len(shape.text.split())
            return slide_count, word_count, True
        else:
            size_kb = os.path.getsize(filepath) / 1024
            if size_kb < 297:
                estimated_slides = max(1, int(size_kb / 30))
                words_per_slide = 20
            elif size_kb < 944:
                estimated_slides = max(1, int(size_kb / 80))
                words_per_slide = 45
            elif size_kb < 3136:
                estimated_slides = max(1, int(size_kb / 120))
                words_per_slide = 65
            else:
                estimated_slides = max(1, int(size_kb / 200))
                words_per_slide = 90
            return estimated_slides, estimated_slides * words_per_slide, False
    except Exception:
        size_kb = os.path.getsize(filepath) / 1024
        estimated_slides = max(1, int(size_kb / 100))
        return estimated_slides, estimated_slides * 45, False

def get_category(ext):
    ext = ext.lower()
    if ext == '.pptx': return "Modern"
    elif ext == '.ppt': return "Legacy"
    elif ext in ['.pps', '.ppsx']: return "Slideshow"
    elif ext in ['.pot', '.potx']: return "Template"
    else: return "Modern"

def create_gauge(confidence):
    """Create a gauge meter chart for confidence."""
    fig, ax = plt.subplots(figsize=(4, 2.5), subplot_kw={'projection': 'polar'})
    fig.patch.set_facecolor('#0a0f0a')
    ax.set_facecolor('#0a0f0a')

    # Gauge goes from 180° to 0° (left to right)
    angles = np.linspace(np.pi, 0, 100)
    
    # Color zones
    ax.barh(1, np.pi * 0.33, left=np.pi * 0.67, height=0.5, color='#ef4444', alpha=0.3)
    ax.barh(1, np.pi * 0.33, left=np.pi * 0.34, height=0.5, color='#f59e0b', alpha=0.3)
    ax.barh(1, np.pi * 0.33, left=np.pi * 0.01, height=0.5, color='#22c55e', alpha=0.3)

    # Needle
    angle = np.pi - (confidence / 100) * np.pi
    ax.annotate('', xy=(angle, 0.9), xytext=(angle, 0),
                arrowprops=dict(arrowstyle='->', color='#22c55e', lw=2.5))

    # Labels
    ax.text(np.pi, 1.6, '0%', color='#4a7a4a', fontsize=8, ha='center', fontfamily='monospace')
    ax.text(np.pi/2, 1.6, '50%', color='#4a7a4a', fontsize=8, ha='center', fontfamily='monospace')
    ax.text(0, 1.6, '100%', color='#4a7a4a', fontsize=8, ha='center', fontfamily='monospace')
    ax.text(np.pi/2, 0.3, f'{confidence:.1f}%', color='#22c55e', fontsize=14,
            ha='center', fontweight='bold', fontfamily='monospace')

    ax.set_ylim(0, 2)
    ax.set_theta_zero_location('E')
    ax.set_theta_direction(-1)
    ax.axis('off')
    plt.tight_layout(pad=0)
    return fig

def create_bar_chart(complexity, density, category, confidence):
    """Create prediction bar chart."""
    fig, axes = plt.subplots(1, 3, figsize=(10, 3))
    fig.patch.set_facecolor('#0a0f0a')

    configs = [
        ('COMPLEXITY', complexity, ['Low', 'Medium', 'High'],
         ['#22c55e', '#f59e0b', '#ef4444']),
        ('DENSITY', density, ['Light', 'Medium', 'Dense'],
         ['#22c55e', '#f59e0b', '#ef4444']),
        ('CATEGORY', category, ['Legacy', 'Modern', 'Slideshow', 'Template'],
         ['#3b82f6', '#22c55e', '#f59e0b', '#a855f7']),
    ]

    for ax, (title, value, options, colors) in zip(axes, configs):
        ax.set_facecolor('#0d150d')
        values = [1 if opt == value else 0.15 for opt in options]
        bar_colors = [c if opt == value else '#1a3a1a' for opt, c in zip(options, colors)]
        bars = ax.bar(options, values, color=bar_colors, edgecolor='#1a3a1a', linewidth=0.5)
        ax.set_title(title, color='#4a7a4a', fontsize=9,
                    fontfamily='monospace', pad=8)
        ax.set_ylim(0, 1.3)
        ax.tick_params(colors='#4a7a4a', labelsize=7)
        for spine in ax.spines.values():
            spine.set_color('#1a3a1a')
        ax.set_facecolor('#0d150d')
        plt.setp(ax.get_xticklabels(), color='#a0c4a0', fontfamily='monospace', fontsize=7)
        plt.setp(ax.get_yticklabels(), color='#4a7a4a', fontsize=6)
        # Highlight selected
        for bar, opt in zip(bars, options):
            if opt == value:
                ax.text(bar.get_x() + bar.get_width()/2, 1.05,
                       '▲', ha='center', color='#22c55e', fontsize=10)

    plt.tight_layout(pad=1.0)
    return fig

def export_pdf(filename, size_mb, slide_count, word_count, avg_words,
               complexity, density, category, confidence, method, proc_time):
    """Generate a PDF report."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_fill_color(10, 15, 10)
    pdf.rect(0, 0, 210, 297, 'F')

    # Title
    pdf.set_font('Courier', 'B', 20)
    pdf.set_text_color(34, 197, 94)
    pdf.cell(0, 15, 'AUTODECK AI - ANALYSIS REPORT', new_x="LMARGIN", new_y="NEXT", align='C')

    pdf.set_font('Courier', '', 10)
    pdf.set_text_color(74, 122, 74)
    pdf.cell(0, 8, 'PRESENTATION INTELLIGENCE SYSTEM | INICAI HACKATHON', ln=True, align='C')
    pdf.ln(5)

    # Divider
    pdf.set_draw_color(34, 197, 94)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(8)

    # File info
    pdf.set_font('Courier', 'B', 11)
    pdf.set_text_color(34, 197, 94)
    pdf.cell(0, 8, 'FILE INTELLIGENCE', ln=True)
    pdf.set_font('Courier', '', 9)
    pdf.set_text_color(160, 196, 160)

    info = [
        ('FILENAME', filename),
        ('FILE SIZE', f'{size_mb} MB'),
        ('SLIDES', str(slide_count)),
        ('WORDS', str(word_count)),
        ('AVG WORDS/SLIDE', f'{avg_words:.1f}'),
        ('EXTRACTION METHOD', method),
        ('PROCESSING TIME', f'{proc_time:.2f}s'),
    ]
    for label, value in info:
        pdf.cell(60, 7, f'{label}:', border=0)
        pdf.cell(0, 7, value, ln=True)

    pdf.ln(5)
    pdf.set_draw_color(34, 197, 94)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(8)

    # Predictions
    pdf.set_font('Courier', 'B', 11)
    pdf.set_text_color(34, 197, 94)
    pdf.cell(0, 8, 'PREDICTIONS', ln=True)
    pdf.set_font('Courier', '', 9)
    pdf.set_text_color(160, 196, 160)

    preds = [
        ('COMPLEXITY', complexity),
        ('DENSITY', density),
        ('CATEGORY', category),
        ('CONFIDENCE', f'{confidence:.1f}%'),
        ('QUALITY SCORE', '0 (No metadata available)'),
    ]
    for label, value in preds:
        pdf.cell(60, 7, f'{label}:', border=0)
        pdf.set_text_color(34, 197, 94)
        pdf.cell(0, 7, value, ln=True)
        pdf.set_text_color(160, 196, 160)

    pdf.ln(5)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(8)

    # Footer
    pdf.set_font('Courier', '', 8)
    pdf.set_text_color(74, 122, 74)
    pdf.cell(0, 6, 'AUTODECK AI | RANDOM FOREST ENGINE | INICAI HACKATHON 2026', ln=True, align='C')
    pdf.cell(0, 6, 'https://huggingface.co/spaces/Drax369/autodeck-ai', ln=True, align='C')

    # Save
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
    pdf.output(tmp.name)
    return tmp.name

def update_history_html(history):
    """Generate HTML for history panel."""
    if not history:
        return """
        <div style='font-family: monospace; color: #2a5a2a; text-align: center; padding: 20px; font-size: 12px;'>
            [ NO ANALYSIS HISTORY ]
        </div>"""

    rows = ""
    for h in reversed(history[-5:]):
        rows += f"""
        <div style='border-bottom: 1px solid #1a3a1a; padding: 8px 0; font-size: 11px;'>
            <div style='color: #22c55e; font-weight: bold;'>{h['filename'][:30]}...</div>
            <div style='color: #4a7a4a; margin-top: 3px;'>
                {h['complexity']} · {h['density']} · {h['category']} · {h['confidence']:.1f}%
            </div>
            <div style='color: #2a4a2a; font-size: 10px;'>{h['time']}</div>
        </div>"""

    return f"""
    <div style='font-family: monospace; background: #0a0f0a; 
                border: 1px solid #1a3a1a; border-radius: 8px; padding: 12px;'>
        <div style='color: #4a7a4a; font-size: 10px; letter-spacing: 0.15em; margin-bottom: 8px;'>
            ◈ RECENT ANALYSES (LAST 5)
        </div>
        {rows}
    </div>"""

def analyze_presentation(file):
    if file is None:
        return (
            """<div style='font-family:monospace; text-align:center; padding:60px 20px;
                color: #2a5a2a; background:#0a0f0a; border:1px solid #1a3a1a;
                border-radius:12px; min-height:200px;'>
                <div style='font-size:28px;'>⬡</div>
                <div style='margin-top:12px; letter-spacing:0.2em;'>SYSTEM STANDBY</div>
                <div style='font-size:11px; margin-top:6px; color:#1a3a1a;'>AWAITING FILE INPUT</div>
            </div>""",
            None, None, None, update_history_html(analysis_history)
        )

    start_time = time.time()

    filepath = file.name
    filename = os.path.basename(filepath)
    ext = os.path.splitext(filename)[1]
    size_kb = os.path.getsize(filepath) / 1024
    size_mb = round(size_kb / 1024, 2)

    slide_count, word_count, exact = extract_features(filepath, ext)
    avg_words = word_count / max(slide_count, 1)

    X = [[size_kb, size_mb, slide_count, word_count, avg_words]]
    complexity = le_complexity.inverse_transform(clf_complexity.predict(X))[0]
    density = le_density.inverse_transform(clf_density.predict(X))[0]
    category = get_category(ext)

    complexity_proba = clf_complexity.predict_proba(X).max()
    density_proba = clf_density.predict_proba(X).max()
    confidence = round((complexity_proba + density_proba) / 2 * 100, 1)

    proc_time = time.time() - start_time
    method = "DIRECT EXTRACTION" if exact else "SIZE-BASED ESTIMATION"

    # Add to history
    import datetime
    analysis_history.append({
        'filename': filename,
        'complexity': complexity,
        'density': density,
        'category': category,
        'confidence': confidence,
        'time': datetime.datetime.now().strftime("%H:%M:%S")
    })

    complexity_colors = {"Low": "#22c55e", "Medium": "#f59e0b", "High": "#ef4444"}
    density_colors = {"Light": "#22c55e", "Medium": "#f59e0b", "Dense": "#ef4444"}
    cc = complexity_colors.get(complexity, "#22c55e")
    dc = density_colors.get(density, "#22c55e")
    bar_color = "#22c55e" if confidence >= 80 else "#f59e0b" if confidence >= 60 else "#ef4444"

    html = f"""
<div style='font-family: "Courier New", monospace; background: #0a0f0a;
            border: 1px solid #1a3a1a; border-radius: 12px; padding: 24px; color: #a0c4a0;'>

    <!-- Header -->
    <div style='border-bottom: 1px solid #1a3a1a; padding-bottom: 14px; margin-bottom: 20px;'>
        <div style='display: flex; justify-content: space-between; align-items: center;'>
            <div>
                <span style='color: #4a7a4a; font-size: 11px; letter-spacing: 0.15em;'>
                    ▶ AUTODECK INTEL SYSTEM v2.0
                </span><br>
                <span style='color: #22c55e; font-size: 16px; font-weight: bold;'>
                    ANALYSIS COMPLETE
                </span>
            </div>
            <div style='text-align: right;'>
                <span style='background: #22c55e; color: #000; font-size: 10px;
                             font-weight: bold; padding: 3px 10px; border-radius: 4px;
                             letter-spacing: 0.1em;'>● ONLINE</span>
                <div style='color: #4a7a4a; font-size: 10px; margin-top: 4px;'>
                    ⏱ {proc_time:.2f}s
                </div>
            </div>
        </div>
    </div>

    <!-- File Intel -->
    <div style='background: #0d150d; border: 1px solid #1a3a1a; border-radius: 8px;
                padding: 16px; margin-bottom: 16px;'>
        <div style='color: #4a7a4a; font-size: 10px; letter-spacing: 0.15em; margin-bottom: 10px;'>
            ◈ FILE INTELLIGENCE
        </div>
        <div style='color: #e0ffe0; font-size: 15px; font-weight: bold;
                    margin-bottom: 12px; word-break: break-all;'>{filename}</div>
        <div style='display: grid; grid-template-columns: repeat(4, 1fr); gap: 12px;'>
            <div style='background: #111a11; border: 1px solid #1a3a1a;
                        border-radius: 6px; padding: 10px; text-align: center;'>
                <div style='color: #4a7a4a; font-size: 9px;'>SIZE</div>
                <div style='color: #22c55e; font-size: 16px; font-weight: bold;'>{size_mb}</div>
                <div style='color: #4a7a4a; font-size: 9px;'>MB</div>
            </div>
            <div style='background: #111a11; border: 1px solid #1a3a1a;
                        border-radius: 6px; padding: 10px; text-align: center;'>
                <div style='color: #4a7a4a; font-size: 9px;'>SLIDES</div>
                <div style='color: #22c55e; font-size: 16px; font-weight: bold;'>{slide_count}</div>
                <div style='color: #4a7a4a; font-size: 9px;'>COUNT</div>
            </div>
            <div style='background: #111a11; border: 1px solid #1a3a1a;
                        border-radius: 6px; padding: 10px; text-align: center;'>
                <div style='color: #4a7a4a; font-size: 9px;'>WORDS</div>
                <div style='color: #22c55e; font-size: 16px; font-weight: bold;'>{word_count}</div>
                <div style='color: #4a7a4a; font-size: 9px;'>TOTAL</div>
            </div>
            <div style='background: #111a11; border: 1px solid #1a3a1a;
                        border-radius: 6px; padding: 10px; text-align: center;'>
                <div style='color: #4a7a4a; font-size: 9px;'>AVG W/S</div>
                <div style='color: #22c55e; font-size: 16px; font-weight: bold;'>{avg_words:.0f}</div>
                <div style='color: #4a7a4a; font-size: 9px;'>PER SLIDE</div>
            </div>
        </div>
        <div style='margin-top: 10px; font-size: 10px; color: #3a6a3a; letter-spacing: 0.1em;'>
            ◈ EXTRACTION: {method} | ⏱ PROCESSED IN {proc_time:.2f}s
        </div>
    </div>

    <!-- Predictions -->
    <div style='display: grid; grid-template-columns: repeat(3, 1fr); gap: 12px; margin-bottom: 16px;'>
        <div style='background: #0d150d; border: 1px solid {cc}44;
                    border-radius: 8px; padding: 16px; text-align: center;
                    box-shadow: 0 0 12px {cc}22;'>
            <div style='color: #4a7a4a; font-size: 9px; letter-spacing: 0.15em; margin-bottom: 8px;'>
                ◈ COMPLEXITY
            </div>
            <div style='font-size: 26px; font-weight: 900; color: {cc};'>{complexity.upper()}</div>
            <div style='margin-top: 8px; height: 3px; background: #1a2e1a; border-radius: 2px;'>
                <div style='height: 100%; border-radius: 2px; background: {cc};
                    width: {"33%" if complexity=="Low" else "66%" if complexity=="Medium" else "100%"};'>
                </div>
            </div>
        </div>
        <div style='background: #0d150d; border: 1px solid {dc}44;
                    border-radius: 8px; padding: 16px; text-align: center;
                    box-shadow: 0 0 12px {dc}22;'>
            <div style='color: #4a7a4a; font-size: 9px; letter-spacing: 0.15em; margin-bottom: 8px;'>
                ◈ DENSITY
            </div>
            <div style='font-size: 26px; font-weight: 900; color: {dc};'>{density.upper()}</div>
            <div style='margin-top: 8px; height: 3px; background: #1a2e1a; border-radius: 2px;'>
                <div style='height: 100%; border-radius: 2px; background: {dc};
                    width: {"33%" if density=="Light" else "66%" if density=="Medium" else "100%"};'>
                </div>
            </div>
        </div>
        <div style='background: #0d150d; border: 1px solid #22c55e44;
                    border-radius: 8px; padding: 16px; text-align: center;
                    box-shadow: 0 0 12px #22c55e22;'>
            <div style='color: #4a7a4a; font-size: 9px; letter-spacing: 0.15em; margin-bottom: 8px;'>
                ◈ CATEGORY
            </div>
            <div style='font-size: 26px; font-weight: 900; color: #22c55e;'>{category.upper()}</div>
            <div style='margin-top: 8px; font-size: 10px; color: #3a6a3a;'>{ext.upper()}</div>
        </div>
    </div>

</div>
"""

    gauge_fig = create_gauge(confidence)
    bar_fig = create_bar_chart(complexity, density, category, confidence)
    history_html = update_history_html(analysis_history)

    return html, gauge_fig, bar_fig, filepath, history_html

def export_report(file):
    if file is None:
        return None
    if not analysis_history:
        return None
    h = analysis_history[-1]
    filepath = file.name
    filename = os.path.basename(filepath)
    ext = os.path.splitext(filename)[1]
    size_kb = os.path.getsize(filepath) / 1024
    size_mb = round(size_kb / 1024, 2)
    slide_count, word_count, exact = extract_features(filepath, ext)
    avg_words = word_count / max(slide_count, 1)
    method = "DIRECT EXTRACTION" if exact else "SIZE-BASED ESTIMATION"

    pdf_path = export_pdf(
        filename, size_mb, slide_count, word_count, avg_words,
        h['complexity'], h['density'], h['category'], h['confidence'],
        method, 0.0
    )
    return pdf_path

CSS = """
.gradio-container {
    background: #060d06 !important;
    max-width: 1100px !important;
    margin: 0 auto !important;
}
footer { display: none !important; }
"""

with gr.Blocks(title="AutoDeck AI") as demo:

    gr.HTML("""
    <div style='font-family: "Courier New", monospace; text-align: center;
                padding: 36px 0 28px; background: #060d06;'>
        <div style='font-size: 11px; color: #3a6a3a; letter-spacing: 0.3em; margin-bottom: 8px;'>
            ▶ INICAI HACKATHON · AUTODECK TRACK · v2.0
        </div>
        <div style='font-size: 36px; font-weight: 900; color: #22c55e;
                    letter-spacing: 0.1em; text-shadow: 0 0 30px #22c55e66;'>
            AUTODECK AI
        </div>
        <div style='font-size: 13px; color: #4a7a4a; letter-spacing: 0.2em; margin-top: 6px;'>
            PRESENTATION INTELLIGENCE SYSTEM
        </div>
        <div style='width: 60px; height: 2px; background: #22c55e;
                    margin: 16px auto 0; border-radius: 2px;'></div>
    </div>
    """)

    with gr.Row():
        # Left panel
        with gr.Column(scale=1):
            gr.HTML("""
            <div style='font-family: "Courier New", monospace; background: #0a0f0a;
                        border: 1px solid #1a3a1a; border-radius: 10px; padding: 16px;
                        margin-bottom: 12px;'>
                <div style='color: #4a7a4a; font-size: 10px; letter-spacing: 0.15em; margin-bottom: 10px;'>
                    ◈ MISSION BRIEF
                </div>
                <div style='color: #a0c4a0; font-size: 12px; line-height: 1.8;'>
                    Upload any PowerPoint file for instant AI-powered analysis.<br><br>
                    The system will classify:<br>
                    <span style='color: #22c55e;'>▸</span> Complexity Level<br>
                    <span style='color: #22c55e;'>▸</span> Content Density<br>
                    <span style='color: #22c55e;'>▸</span> File Category<br>
                    <span style='color: #22c55e;'>▸</span> Confidence Score
                </div>
            </div>
            """)

            file_input = gr.File(
                label="◈ UPLOAD TARGET FILE",
                file_types=[".ppt", ".pptx", ".pps", ".ppsx", ".pot", ".pptm"]
            )
            analyze_btn = gr.Button("▶ INITIALIZE ANALYSIS", variant="primary", size="lg")
            export_btn = gr.Button("🖨️ EXPORT PDF REPORT", variant="secondary", size="sm")
            pdf_output = gr.File(label="📄 Download Report", visible=True)

            gr.HTML("""
            <div style='font-family: "Courier New", monospace; font-size: 10px;
                        color: #3a5a3a; text-align: center; margin-top: 8px; letter-spacing: 0.1em;'>
                SUPPORTED: .PPT .PPTX .PPS .PPSX .POT .PPTM
            </div>
            """)

            # History panel
            history_html = gr.HTML(update_history_html([]))

        # Right panel
        with gr.Column(scale=2):
            result_html = gr.HTML("""
            <div style='font-family: "Courier New", monospace; background: #0a0f0a;
                        border: 1px solid #1a3a1a; border-radius: 12px; padding: 40px 20px;
                        text-align: center; color: #3a6a3a; min-height: 200px;'>
                <div style='font-size: 32px; margin-bottom: 16px;'>⬡</div>
                <div style='font-size: 13px; letter-spacing: 0.2em;'>SYSTEM STANDBY</div>
                <div style='font-size: 11px; color: #1a3a1a; margin-top: 8px;'>AWAITING FILE INPUT</div>
            </div>
            """)

            with gr.Row():
                gauge_plot = gr.Plot(label="◈ CONFIDENCE GAUGE")
                bar_plot = gr.Plot(label="◈ PREDICTION BREAKDOWN")

    file_store = gr.State()

    analyze_btn.click(
        fn=analyze_presentation,
        inputs=[file_input],
        outputs=[result_html, gauge_plot, bar_plot, file_store, history_html]
    )

    export_btn.click(
        fn=export_report,
        inputs=[file_input],
        outputs=[pdf_output]
    )

    gr.HTML("""
    <div style='font-family: "Courier New", monospace; text-align: center;
                padding: 20px 0 8px; font-size: 10px; color: #2a4a2a;
                letter-spacing: 0.15em; border-top: 1px solid #0d1a0d; margin-top: 20px;'>
        AUTODECK AI v2.0 · RANDOM FOREST ENGINE · INICAI HACKATHON 2026 · ALL SYSTEMS NOMINAL
    </div>
    """)

demo.launch(css=CSS)